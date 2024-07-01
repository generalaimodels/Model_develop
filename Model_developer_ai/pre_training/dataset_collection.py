
import os
import csv
import pandas as pd
import fitz
import json
import yaml
import re
from concurrent.futures import ThreadPoolExecutor
import unicodedata
from typing import List
from pptx import Presentation
from datasets import load_dataset
from typing import Dict, List ,Optional
from collections import defaultdict
from pathlib import Path

def read_pptx(file):
    """Custom function to read .pptx files with python-pptx"""
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


EXTENSION_READERS = {
    '.ipynb': lambda f: json.load(f),
    '.md': lambda f: f.read(),
    '.py': lambda f: f.read(),
    '.csv': lambda f: pd.read_csv(f),
    '.json': lambda f: json.load(f),
    '.yaml': lambda f: yaml.safe_load(f),
    '.txt': lambda f: f.read(),
    '.xml': lambda f: f.read(),
    '.html': lambda f: f.read(),
    '.css': lambda f: f.read(),
    '.js': lambda f: f.read(),
    '.java': lambda f: f.read(),
    '.cpp': lambda f: f.read(),
    '.h': lambda f: f.read(),
    '.php': lambda f: f.read(),
    '.rb': lambda f: f.read(),
    '.sql': lambda f: f.read(),
    '.xls': lambda f: pd.read_excel(f),
    '.xlsx': lambda f: pd.read_excel(f),
    '.html': lambda f: f.read(),
    '.css': lambda f: f.read(),
    '.js': lambda f: f.read(),
    '.jsp': lambda f: f.read(),
    '.jspx': lambda f: f.read(),
    '.vue': lambda f: f.read(),
    '.ejs': lambda f: f.read(),
    '.erb': lambda f: f.read(),
    '.aspx': lambda f: f.read(),
    '.jspx': lambda f: f.read(),
    '.php': lambda f: f.read(),
    '.java': lambda f: f.read(),
    '.cpp': lambda f: f.read(),
    '.h': lambda f: f.read(),
    '.c': lambda f: f.read(),
    '.cc': lambda f: f.read(),
    '.cxx': lambda f: f.read(),
    '.m': lambda f: f.read(),
    '.mm': lambda f: f.read(),
    '.swift': lambda f: f.read(),
    '.objcpp': lambda f: f.read(),
    '.cs': lambda f: f.read(),
    '.go': lambda f: f.read(),
    '.rs': lambda f: f.read(),
    '.kt': lambda f: f.read(),
    '.dart': lambda f: f.read(),
    '.rb': lambda f: f.read(),
    '.pl': lambda f: f.read(),
    '.pm': lambda f: f.read(),
    '.py': lambda f: f.read(),
    '.sh': lambda f: f.read(),
    '.sql': lambda f: f.read(),
    '.sqlite': lambda f: f.read(),
    '.xml': lambda f: f.read(),
    '.ppt': lambda f: read_pptx(f),
    '.pptx': lambda f: read_pptx(f)
}

EXTENSION_PATTERN = r".*\.(md|py|csv|json|yaml|txt|xml|html|css|js|java|cpp|h|php|rb|sql|xls|xlsx|ppt|pptx|ipynb|jsp|jspx|vue|ejs|erb|aspx|c|cc|cxx|m|mm|swift|objcpp|cs|go|rs|kt|dart|pl|pm|sh|sqlite)$"

def get_files_with_extensions(dir_path: str) -> Dict[str, List[str]]:
    ext_files = defaultdict(list)
    for root, dirs, files in os.walk(dir_path):
        for file in files:
            file_path = os.path.join(root, file).replace("\\", "/")
            _, ext = os.path.splitext(file)
            ext_files[ext].append(file_path)
    return ext_files


def write_to_csv(file_path: str, ext_files: Dict[str, List[str]]) -> None:
    with open(file_path, 'w', encoding='utf-8',newline='') as csvfile:
        writer = csv.writer(csvfile)
        max_len = max(len(v) for v in ext_files.values())
        writer.writerow(ext_files.keys())
        for i in range(max_len):
            row = [ext_files[k][i] if i < len(ext_files[k]) else '' for k in ext_files.keys()]
            writer.writerow(row)




# import re
# import string
# from nltk.tokenize import word_tokenize
# from nltk.corpus import stopwords
# from nltk.stem import WordNetLemmatizer
# from bs4 import BeautifulSoup
# import html

# def clean_text(text: str) -> str:
#     """
#     Clean the extracted text from the Document.
#     This function can be customized based on the cleaning requirements.

#     Parameters:
#     - text (str): The text extracted from the PDF.

#     Returns:
#     - str: The cleaned text.
#     """
#     # Remove HTML tags
#     cleaned_text = BeautifulSoup(text, 'html.parser').get_text()

#     # Remove leading/trailing whitespace
#     cleaned_text = cleaned_text.strip()

#     # Replace multiple consecutive whitespace characters with a single space
#     cleaned_text = ' '.join(cleaned_text.split())

#     # Remove non-printable characters
#     cleaned_text = ''.join(char for char in cleaned_text if char.isprintable())

#     # Remove any remaining control characters
#     cleaned_text = ''.join(char for char in cleaned_text if not char.isspace() or char == ' ')

#     # Remove punctuation
#     cleaned_text = cleaned_text.translate(str.maketrans('', '', string.punctuation))

#     # Tokenize the text
#     tokens = word_tokenize(cleaned_text)

#     # Remove stopwords
#     stop_words = set(stopwords.words('english'))
#     tokens = [word for word in tokens if word not in stop_words]

#     # Lemmatize the tokens
#     lemmatizer = WordNetLemmatizer()
#     tokens = [lemmatizer.lemmatize(word) for word in tokens]

#     # Join the tokens back into a string
#     cleaned_text = ' '.join(tokens)

#     # Check if the cleaned text is empty
#     if not cleaned_text:
#         return ''

#     # Capitalize the first character of the cleaned text
#     cleaned_text = cleaned_text[0].upper() + cleaned_text[1:]

#     # Check if the cleaned text ends with proper punctuation
#     if not cleaned_text.endswith(('.', '!', '?')):
#         cleaned_text += '.'

#     return cleaned_text

def clean_text(text: str) -> Optional[str]:
    """
    Clean the extracted text from the Document using advanced techniques.
    This function can be customized based on the cleaning requirements.

    Parameters:
    - text (str): The text extracted from the PDF.

    Returns:
    - Optional[str]: The cleaned text, or None if the input is invalid.
    """
    if not isinstance(text, str):
        return None

    # Remove leading/trailing whitespace
    cleaned_text = text.strip()

    # # Remove non-printable characters and control characters using regex
    # cleaned_text = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', cleaned_text)

    # Normalize Unicode characters
    # cleaned_text = unicodedata.normalize('NFKC', cleaned_text)

    # Remove extra whitespace and line breaks using regex
    # cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

    # # Remove URLs using regex
    # cleaned_text = re.sub(r'https?://\S+|www\.\S+', '', cleaned_text)

    # # Remove email addresses using regex
    # cleaned_text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '', cleaned_text)

    # # Remove phone numbers using regex
    # cleaned_text = re.sub(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '', cleaned_text)

    # # Remove special characters and punctuation using regex
    # cleaned_text = re.sub(r'[^a-zA-Z0-9\s]', '', cleaned_text)

    # Convert to lowercase
    cleaned_text = cleaned_text.lower()

    # Remove stopwords (optional)
    # stopwords = ['the', 'and', 'in', 'on', 'at', 'to', 'of', 'for', 'with', 'by', 'from', 'up', 'about', 'into', 'over', 'after', 'beneath', 'under', 'above', 'across', 'before', 'behind', 'below', 'beside', 'between', 'beyond', 'near', 'toward', 'through', 'during', 'except', 'inside', 'outside', 'since', 'until', 'while', 'within', 'without', 'along', 'among', 'against', 'around', 'despite', 'toward', 'upon', 'considering', 'following', 'including', 'regarding', 'according', 'aside', 'away', 'because', 'besides', 'concerning', 'due', 'like', 'next', 'off', 'onto', 'out', 'throughout', 'underneath', 'unlike', 'until', 'upon', 'versus', 'via', 'whether', 'though', 'although', 'even', 'whereas', 'while', 'whilst']
    # cleaned_text = ' '.join(word for word in cleaned_text.split() if word.lower() not in stopwords)

    # Check if the cleaned text is empty
    if not cleaned_text:
        return None

    return cleaned_text
def split_and_save_text(cleaned_text: str, base_output_path: Path, max_size_bytes: int = 50 * 1024 * 1024) -> None:
    """
    Split the cleaned text into multiple files, each smaller than the specified max size, and save them.

    Parameters:
    - cleaned_text (str): The cleaned text to be split and saved.
    - base_output_path (Path): The base path where the text files will be saved.
    - max_size_bytes (int): Maximum size of the text file in bytes.
    """
    part_num = 1
    text_part = ""
    for line in cleaned_text.split('\n'):
        if len(text_part.encode('utf-8')) + len(line.encode('utf-8')) < max_size_bytes:
            text_part += line + '\n'
        else:
            # Save the current part and start a new one
            output_path = base_output_path.with_suffix(f'.part{part_num}.txt')
            with open(output_path, 'w', encoding='utf-8') as file:
                file.write(text_part)
            part_num += 1
            text_part = line + '\n' # Start new part with the current line

    # Save the last part
    if text_part:
        output_path = base_output_path.with_suffix(f'.part{part_num}.txt')
        with open(output_path, 'w', encoding='utf-8') as file:
            file.write(text_part)


def convert_pdf_to_text(pdf_path: str, output_folder: str) -> None:
    """
    Convert a PDF file to text files, splitting contents to ensure each resulting file is less than 50 MB.

    Parameters:
    - pdf_path (str): Path to the PDF file.
    - output_folder (str): Path to the folder where the text files will be saved.
    """
    # Ensure the output folder exists, create it if it does not
    output_folder_path = Path(output_folder)
    output_folder_path.mkdir(parents=True, exist_ok=True)

    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()

        cleaned_text = clean_text(text)

        base_output_path = output_folder_path / Path(pdf_path).stem
        split_and_save_text(cleaned_text, base_output_path)
    except Exception as e:
        # print(f"An error occurred while converting {pdf_path}: {str(e)}")
        return None


def process_pdfs_from_csv(csv_path: str, output_folder: str) -> None:
    """
    Process PDFs listed in a CSV file, converting them to text files and ensuring each part is less than 50 MB.

    Parameters:
    - csv_path (str): Path to the CSV file containing paths to PDF files.
    - output_folder (str): Path to the folder where text files will be stored.
    """
    encodings = ['utf-8', 'latin1', 'utf-16', 'utf-32','ascii']

    for encoding in encodings:
        try:
            pdf_paths = pd.read_csv(csv_path, encoding=encoding)
            break  # If successful, break out of the loop
        except UnicodeDecodeError:
            continue  # If the encoding fails, try the next one
    else:
        raise ValueError("Unable to decode the CSV file with the provided encodings.")

    for pdf_path in pdf_paths['.pdf']:
        convert_pdf_to_text(pdf_path, output_folder)

def list_files_with_extensions(directory_path):
    try:
        files = os.listdir(directory_path)
        return [file for file in files if re.match(EXTENSION_PATTERN, file)]
    except FileNotFoundError:
        # print(f"The directory {directory_path} was not found.")
        return None

def read_file_content(directory_path, filename):
    try:
        extension = os.path.splitext(filename)[1]
        with open(os.path.join(directory_path, filename), 'r',encoding='utf-8') as file:
            file_reader = EXTENSION_READERS.get(extension)
            return file_reader(file) if file_reader else None
    except Exception as e:
        # print(f"An error occurred while reading the file {filename}: {e}")
        return None



def process_files_txtfile(directory_path: str, user_folder: str, txt_file_counter: int = 1, txt_file_size: int = 0) -> Optional[str]:
    """
    This function recursively processes all files in a given directory and its subdirectories,
    and writes their content to a user-specific text file. Each text file is ensured to be less than 50 MB in size.

    Args:
    directory_path (str): The path to the directory containing the files to be processed.
    user_folder (str): The name of the user-specific folder where the text files will be written.
    txt_file_counter (int): The counter for the current text file.
    txt_file_size (int): The current size of the text file.

    Returns:
    str: The path to the user-specific folder, or None if an error occurred.
    """
    
    user_folder_path = os.path.join(directory_path, user_folder)
    os.makedirs(user_folder_path, exist_ok=True)
    files = list_files_with_extensions(directory_path)

    if files is None:
        return

    for filename in files:
        content = read_file_content(directory_path, filename)
        if content is not None:
            
            if txt_file_size >= 5*1024:
                txt_file_counter += 1
                txt_file_size = 0
            txt_file_path = os.path.join(user_folder_path, f"{txt_file_counter}.txt")
            with open(txt_file_path, "a", encoding='utf-8') as f:
                if isinstance(content, pd.DataFrame):
                    content_csv = content.to_csv(index=False)
                    f.write(content_csv)
                    txt_file_size += len(content_csv)
                elif isinstance(content, dict):
                    content_str = json.dumps(content, indent=4)
                    f.write(content_str)
                    txt_file_size += len(content_str)
                else:
                    content_str = str(content)
                    f.write(content_str)
                    txt_file_size += len(content_str)
    for subdir in [d for d in os.listdir(directory_path) if os.path.isdir(os.path.join(directory_path, d))]:
        subdir_path = os.path.join(directory_path, subdir)
        txt_file_counter, txt_file_size = process_files_txtfile(subdir_path, user_folder, txt_file_counter, txt_file_size)
    return txt_file_counter, txt_file_size





def reformat_txt_file(file_path: str, new_folder: str) -> None:
    """
    Reformats a single .txt file to give it a clean, professional look.
    Saves the reformatted file in a new folder while keeping the original file intact.

    Args:
        file_path (str): Path to the .txt file to reformat.
        new_folder (str): Path to the folder where the reformatted file will be saved.

    Returns:
        None
    """
    try:
        with open(file_path, "r", encoding='utf-8') as file:
            content = file.read()
        content = content.strip()
        content = " ".join(content.split())
        content = content.replace("\n", "\n\n")
        lines = content.split("\n")
        if len(lines) > 50:
            content = ""
            for i, line in enumerate(lines, start=1):
                content += line + "\n"
                if i % 50 == 0:
                    content += f"\nPage {i // 50}\n\n"
        new_filename = "_".join(os.path.basename(file_path).split())
        new_file_path = os.path.join(new_folder, new_filename)
        with open(new_file_path, "w", encoding='utf-8') as file:
            file.write(content)
    except Exception as e:
        print(f"Error processing file: {file_path}. Error: {str(e)}")

def reformat_txt_files(folders: List[str], max_workers: int = 4) -> None:
    """
    Reformats .txt files in the specified folders to give them a clean, professional look.
    Saves the reformatted files in new folders while keeping the original files intact.

    Args:
        folders (List[str]): List of folder paths containing the .txt files to reformat.
        max_workers (int): Maximum number of worker threads to use for concurrent processing.

    Returns:
        None
    """
    for folder in folders:
        new_folder = os.path.join(folder, "reformatted")
        os.makedirs(new_folder, exist_ok=True)

        txt_files = [os.path.join(folder, filename) for filename in os.listdir(folder) if filename.endswith(".txt")]

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = []
            for file_path in txt_files:
                future = executor.submit(reformat_txt_file, file_path, new_folder)
                futures.append(future)

            for future in futures:
                future.result()


def loading_folder_using_datasets(folder_path:str):
    dataset = load_dataset('text', data_files=folder_path+'/*.txt',data_dir="./output_dataset")
    return dataset





# # #===============================================|$$ development is going on $$ |===========================
# DIR_PATH="C:/Users/heman/Desktop/Coding/dataset"
# DIR_OUTPUT="C:/Users/heman/Desktop/Coding/dataset/output"
# CSV_FILE_PATH="C:/Users/heman/Desktop/Coding/dataset/csv_file.csv "
# FOLDER_TO_PROCESS=[DIR_OUTPUT]
# ext_files = get_files_with_extensions(DIR_PATH)
# write_to_csv(CSV_FILE_PATH, ext_files)
# process_pdfs_from_csv(csv_path=CSV_FILE_PATH, output_folder=DIR_OUTPUT)
# process_files_txtfile(DIR_PATH,  DIR_OUTPUT)
# reformat_txt_files(FOLDER_TO_PROCESS)
# dataset=loading_folder_using_datasets(folder_path=f'{DIR_OUTPUT}/reformatted')
# print(dataset)