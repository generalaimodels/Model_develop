import os
import csv
import pandas as pd
import fitz
import json
import yaml
import torch
import urllib.request
import re
import multiprocessing
from typing import List
from urllib.parse import quote_plus
from pptx import Presentation
from io import StringIO
from datasets import load_dataset,Dataset
from tqdm.notebook import tqdm
from typing import Dict, List ,Optional
from collections import defaultdict
from pathlib import Path
from typing import Optional, List, Tuple
from langchain_community.document_loaders import DirectoryLoader
from langchain.text_splitter import CharacterTextSplitter
from langchain.text_splitter import RecursiveCharacterTextSplitter
from transformers import AutoTokenizer
from langchain.vectorstores.faiss import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores.utils import DistanceStrategy
from transformers import pipeline
from transformers import AutoTokenizer, AutoModelForCausalLM, BitsAndBytesConfig
from langchain.docstore.document import Document as LangchainDocument
from datetime import date, timedelta


def read_pptx(file):
    """Custom function to read .pptx files with python-pptx"""
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)
EMBEDDING_MODEL_NAME = "thenlper/gte-small"
MARKDOWN_SEPARATORS = [
    "\n#{1,6} ",
    "```\n",
    "\n\\*\\*\\*+\n",
    "\n---+\n",
    "\n___+\n",
    "\n\n",
    "\n",
    " ",
    "",
]
docs_processed=[]

READER_MODEL_NAME = "HuggingFaceH4/zephyr-7b-beta"
EXTENSION_READERS = {
    '.ipynb': lambda f: json.load(f),
    '.md': lambda f: f.read(),
    '.py': lambda f: f.read(),
    '.csv': lambda f: pd.read_csv(f),
    '.json': lambda f: json.load(f),
    '.yaml': lambda f: yaml.safe_load(f),
    '.txt': lambda f: f.read(),
    '.xml': lambda f: f.read(),
    '.html': lambda f: f.read(),
    '.css': lambda f: f.read(),
    '.js': lambda f: f.read(),
    '.java': lambda f: f.read(),
    '.cpp': lambda f: f.read(),
    '.h': lambda f: f.read(),
    '.php': lambda f: f.read(),
    '.rb': lambda f: f.read(),
    '.sql': lambda f: f.read(),
    '.xls': lambda f: pd.read_excel(f),
    '.xlsx': lambda f: pd.read_excel(f),
    '.html': lambda f: f.read(),
    '.css': lambda f: f.read(),
    '.js': lambda f: f.read(),
    '.jsp': lambda f: f.read(),
    '.jspx': lambda f: f.read(),
    '.vue': lambda f: f.read(),
    '.ejs': lambda f: f.read(),
    '.erb': lambda f: f.read(),
    '.aspx': lambda f: f.read(),
    '.jspx': lambda f: f.read(),
    '.php': lambda f: f.read(),
    '.java': lambda f: f.read(),
    '.cpp': lambda f: f.read(),
    '.h': lambda f: f.read(),
    '.c': lambda f: f.read(),
    '.cc': lambda f: f.read(),
    '.cxx': lambda f: f.read(),
    '.m': lambda f: f.read(),
    '.mm': lambda f: f.read(),
    '.swift': lambda f: f.read(),
    '.objcpp': lambda f: f.read(),
    '.cs': lambda f: f.read(),
    '.go': lambda f: f.read(),
    '.rs': lambda f: f.read(),
    '.kt': lambda f: f.read(),
    '.dart': lambda f: f.read(),
    '.rb': lambda f: f.read(),
    '.pl': lambda f: f.read(),
    '.pm': lambda f: f.read(),
    '.py': lambda f: f.read(),
    '.sh': lambda f: f.read(),
    '.sql': lambda f: f.read(),
    '.sqlite': lambda f: f.read(),
    '.xml': lambda f: f.read(),
    '.ppt': lambda f: read_pptx(f),
    '.pptx': lambda f: read_pptx(f)
}

# Utilize regular expressions to match any of the file extensions
EXTENSION_PATTERN = r".*\.(md|py|csv|json|yaml|txt|xml|html|css|js|java|cpp|h|php|rb|sql|xls|xlsx|ppt|pptx|ipynb|jsp|jspx|vue|ejs|erb|aspx|c|cc|cxx|m|mm|swift|objcpp|cs|go|rs|kt|dart|pl|pm|sh|sqlite)$"

def get_files_with_extensions(dir_path: str) -> Dict[str, List[str]]:
    ext_files = defaultdict(list)
    for root, dirs, files in os.walk(dir_path):
        for file in files:
            file_path = os.path.join(root, file).replace("\\", "/")
            _, ext = os.path.splitext(file)
            ext_files[ext].append(file_path)
    return ext_files


def write_to_csv(file_path: str, ext_files: Dict[str, List[str]]) -> None:
    with open(file_path, 'w', newline='') as csvfile:
        writer = csv.writer(csvfile)
        max_len = max(len(v) for v in ext_files.values())
        writer.writerow(ext_files.keys())
        for i in range(max_len):
            row = [ext_files[k][i] if i < len(ext_files[k]) else '' for k in ext_files.keys()]
            writer.writerow(row)
            


def clean_text(text: str) -> str:
    """
    Clean the extracted text from the PDF.
    This function can be customized based on the cleaning requirements.

    Parameters:
    - text (str): The text extracted from the PDF.

    Returns:
    - str: The cleaned text.
    """
    cleaned_text = ' '.join(text.split())  # Removing extra whitespaces
    # Add more cleaning rules as needed.
    return cleaned_text

def split_and_save_text(cleaned_text: str, base_output_path: Path, max_size_bytes: int = 50 * 1024 * 1024) -> None:
    """
    Split the cleaned text into multiple files, each smaller than the specified max size, and save them.

    Parameters:
    - cleaned_text (str): The cleaned text to be split and saved.
    - base_output_path (Path): The base path where the text files will be saved.
    - max_size_bytes (int): Maximum size of the text file in bytes.
    """
    part_num = 1
    text_part = ""
    for line in cleaned_text.split('\n'):
        if len(text_part.encode('utf-8')) + len(line.encode('utf-8')) < max_size_bytes:
            text_part += line + '\n'
        else:
            # Save the current part and start a new one
            output_path = base_output_path.with_suffix(f'.part{part_num}.txt')
            with open(output_path, 'w', encoding='utf-8') as file:
                file.write(text_part)
            part_num += 1
            text_part = line + '\n' # Start new part with the current line

    # Save the last part
    if text_part:
        output_path = base_output_path.with_suffix(f'.part{part_num}.txt')
        with open(output_path, 'w', encoding='utf-8') as file:
            file.write(text_part)


def convert_pdf_to_text(pdf_path: str, output_folder: str) -> None:
    """
    Convert a PDF file to text files, splitting contents to ensure each resulting file is less than 50 MB.

    Parameters:
    - pdf_path (str): Path to the PDF file.
    - output_folder (str): Path to the folder where the text files will be saved.
    """
    # Ensure the output folder exists, create it if it does not
    output_folder_path = Path(output_folder)
    # output_folder_path.mkdir(parents=True, exist_ok=True)

    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()

        cleaned_text = clean_text(text)

        base_output_path = output_folder_path / Path(pdf_path).stem
        split_and_save_text(cleaned_text, base_output_path)
    except Exception as e:
        print(f"An error occurred while converting {pdf_path}: {str(e)}")


def process_pdfs_from_csv(csv_path: str, output_folder: str) -> None:
    """
    Process PDFs listed in a CSV file, converting them to text files and ensuring each part is less than 50 MB.

    Parameters:
    - csv_path (str): Path to the CSV file containing paths to PDF files.
    - output_folder (str): Path to the folder where text files will be stored.
    """
    encodings = ['utf-8', 'latin1', 'utf-16', 'utf-32','ascii'] 
    
    for encoding in encodings:
        try:
            pdf_paths = pd.read_csv(csv_path, encoding=encoding)
            break  # If successful, break out of the loop
        except UnicodeDecodeError:
            continue  # If the encoding fails, try the next one
    else:
        raise ValueError("Unable to decode the CSV file with the provided encodings.")
    
    for pdf_path in pdf_paths['.pdf']:
        convert_pdf_to_text(pdf_path, output_folder)

def list_files_with_extensions(directory_path):
    try:
        files = os.listdir(directory_path)
        return [file for file in files if re.match(EXTENSION_PATTERN, file)]
    except FileNotFoundError:
        print(f"The directory {directory_path} was not found.")
        return None

def read_file_content(directory_path, filename):
    try:
        extension = os.path.splitext(filename)[1]
        with open(os.path.join(directory_path, filename), 'r',encoding='utf-8') as file:
            file_reader = EXTENSION_READERS.get(extension)
            return file_reader(file) if file_reader else None
    except Exception as e:
        print(f"An error occurred while reading the file {filename}: {e}")



def process_files_txtfile(directory_path: str, user_folder: str, txt_file_counter: int = 1, txt_file_size: int = 0) -> Optional[str]:
    """
    This function recursively processes all files in a given directory and its subdirectories,
    and writes their content to a user-specific text file. Each text file is ensured to be less than 50 MB in size.

    Args:
    directory_path (str): The path to the directory containing the files to be processed.
    user_folder (str): The name of the user-specific folder where the text files will be written.
    txt_file_counter (int): The counter for the current text file.
    txt_file_size (int): The current size of the text file.

    Returns:
    str: The path to the user-specific folder, or None if an error occurred.
    """
    # Create the user-specific folder if it doesn't exist
    user_folder_path = os.path.join(directory_path, user_folder)
    os.makedirs(user_folder_path, exist_ok=True)

    # Get a list of all files in the directory
    files = list_files_with_extensions(directory_path)

    if files is None:
        return

    for filename in files:
        content = read_file_content(directory_path, filename)
        if content is not None:
            # Create a new text file if the size is over 50 MB
            if txt_file_size >= 50 * 1024 * 1024:
                txt_file_counter += 1
                txt_file_size = 0

            # Open the text file in append mode
            txt_file_path = os.path.join(user_folder_path, f"{txt_file_counter}.txt")
            with open(txt_file_path, "a") as f:
                if isinstance(content, pd.DataFrame):
                    # Convert DataFrame to CSV string without index and write to file
                    content_csv = content.to_csv(index=False)
                    f.write(content_csv)
                    txt_file_size += len(content_csv)
                elif isinstance(content, dict):
                    # Convert dict to pretty-printed string using json.dumps and write to file
                    content_str = json.dumps(content, indent=4)
                    f.write(content_str)
                    txt_file_size += len(content_str)
                else:
                    # If it's not a DataFrame or dict, convert it to a string
                    content_str = str(content)
                    f.write(content_str)
                    txt_file_size += len(content_str)

            print(f"--- File: {filename} ---")
            print(f"Content written to {txt_file_path}")
            print("-------------------------------\n")

    # Recursively process subdirectories
    for subdir in [d for d in os.listdir(directory_path) if os.path.isdir(os.path.join(directory_path, d))]:
        subdir_path = os.path.join(directory_path, subdir)
        txt_file_counter, txt_file_size = process_files_txtfile(subdir_path, user_folder, txt_file_counter, txt_file_size)

    return txt_file_counter, txt_file_size
def loading_folder_using_datasets(folder_path:str):


    dataset = load_dataset('text', data_files=folder_path+'/*.txt')
    return dataset

import os
from typing import List

def reformat_txt_files(folders: List[str]) -> None:
    """
    Reformats .txt files in the specified folders to give them a clean, professional look.

    Args:
        folders (List[str]): List of folder paths containing the .txt files to reformat.

    Returns:
        None
    """
    for folder in folders:
        for filename in os.listdir(folder):
            if filename.endswith(".txt"):
                file_path = os.path.join(folder, filename)
                with open(file_path, "r") as file:
                    content = file.read()

                # Clean up the content
                content = content.strip()  # Remove leading/trailing whitespace
                content = " ".join(content.split())  # Remove extra whitespace between words
                content = content.replace("\n", "\n\n")  # Add blank line between paragraphs

                # Add page numbers if the document is multiple pages
                lines = content.split("\n")
                if len(lines) > 50:
                    content = ""
                    for i, line in enumerate(lines, start=1):
                        content += line + "\n"
                        if i % 50 == 0:
                            content += f"\nPage {i // 50}\n\n"

                # Save the reformatted content
                new_filename = "_".join(filename.split())  # Replace spaces with underscores
                new_file_path = os.path.join(folder, new_filename)
                with open(new_file_path, "w") as file:
                    file.write(content)

                print(f"Reformatted: {new_file_path}")




def split_documents(
    chunk_size: int,
    knowledge_base,
    tokenizer_name: Optional[str] = EMBEDDING_MODEL_NAME,
) -> List:
    """
    This function splits documents into chunks of maximum size `chunk_size` tokens and return a list of documents.

    Parameters:
    chunk_size (int): Maximum size of each chunk
    knowledge_base (List): List of documents to be processed
    tokenizer_name (str, optional): Name of the tokenizer. Defaults to EMBEDDING_MODEL_NAME.

    Returns:
    List: List of processed documents
    """
    text_splitter = RecursiveCharacterTextSplitter.from_huggingface_tokenizer(
        AutoTokenizer.from_pretrained(tokenizer_name),
        chunk_size=chunk_size,
        chunk_overlap=int(chunk_size / 10),
        add_start_index=True,
        strip_whitespace=True,
        separators=MARKDOWN_SEPARATORS,
    )

    docs_processed = []
    for doc in knowledge_base:
        docs_processed += text_splitter.split_documents([doc])
    print("****completed*** all files files")


    # Remove duplicates
    unique_texts = {}
    docs_processed_unique = []
    for doc in docs_processed:
        if doc.page_content not in unique_texts:
            unique_texts[doc.page_content] = True
            docs_processed_unique.append(doc)

    return docs_processed_unique


# Example usage
folders_to_process = [
    "/path/to/folder1",
    "/path/to/folder2",
    "/path/to/folder3"
]
reformat_txt_files(folders_to_process)


embedding_model = HuggingFaceEmbeddings(
    model_name=EMBEDDING_MODEL_NAME,
    multi_process=True,
    model_kwargs={"device": "cuda"},
    encode_kwargs={"normalize_embeddings": True},  # set True for cosine similarity
)
bnb_config = BitsAndBytesConfig(
    load_in_4bit=True,
    bnb_4bit_use_double_quant=True,
    bnb_4bit_quant_type="nf4",
    bnb_4bit_compute_dtype=torch.bfloat16,
)
model = AutoModelForCausalLM.from_pretrained(READER_MODEL_NAME, quantization_config=bnb_config)
tokenizer = AutoTokenizer.from_pretrained(READER_MODEL_NAME)
READER_LLM = pipeline(
    model=model,
    tokenizer=tokenizer,
    task="text-generation",
    do_sample=True,
    temperature=0.1,
    repetition_penalty=1.1,
    return_full_text=False,
    max_new_tokens=1024,
)

prompt_in_chat_format = [
    {
        "role": "system",
        "content": """Using the information contained in the context,
give a comprehensive answer to the question.
Respond only to the question asked, response should be concise and relevant to the question.
Provide the number of the source document when relevant.
If the answer cannot be deduced from the context, do not give an answer.""",
    },
    {
        "role": "user",
        "content": """Context:
{context}
---
Now here is the question you need to answer.

Question: {question}""",
    },
]

dir_path ="/home/hemanth/Hemanth"
dir_output='/home/hemanth/Hemanth/data' # replace with your directory path
csv_file_path = '/home/hemanth/Hemanth/csvfile.csv'  # replace with your CSV file path
ext_files = get_files_with_extensions(dir_path)
write_to_csv(csv_file_path, ext_files)
process_pdfs_from_csv(csv_path=csv_file_path, output_folder=dir_output)
process_files_txtfile(dir_path,  dir_output)
dataset=loading_folder_using_datasets(folder_path=dir_output)

knowledge_base = [
LangchainDocument(page_content=doc["text"])
for doc in tqdm(dataset['train'])
]
docs_processed = split_documents(
1024, 
knowledge_base,
tokenizer_name=EMBEDDING_MODEL_NAME,
)
    
    
KNOWLEDGE_VECTOR_DATABASE = FAISS.from_documents(
    docs_processed, embedding_model, distance_strategy=DistanceStrategy.COSINE)




   
        




if __name__ == "__main__":
    multiprocessing.freeze_support()
    user_query=" joking Engineer"
    retrieved_docs = KNOWLEDGE_VECTOR_DATABASE.similarity_search(query=user_query, k=5)
    RAG_PROMPT_TEMPLATE = tokenizer.apply_chat_template(
    prompt_in_chat_format, tokenize=False, add_generation_prompt=True
    )
    retrieved_docs_text = [
    doc.page_content for doc in retrieved_docs
    ]
    context = "\nExtracted documents:\n"
    context += "".join([f"Document {str(i)}:::\n" + doc for i, doc in enumerate(retrieved_docs_text)])
    final_prompt = RAG_PROMPT_TEMPLATE.format(
        question=user_query, context=context
    )

    # Redact an answer
    answer = READER_LLM(final_prompt)[0]["generated_text"]
    print(answer)
