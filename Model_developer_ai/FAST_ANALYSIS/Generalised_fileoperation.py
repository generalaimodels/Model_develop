import os
import sys
import shutil
import subprocess
import logging
import csv
import pandas as pd
import fitz
import json
import yaml
import re
import unicodedata
from typing import List
from pptx import Presentation
from datasets import load_dataset
from typing import Dict, List ,Optional
from collections import defaultdict
from pathlib import Path
from typing import List, Dict,Optional
from datetime import datetime
from datasets import load_dataset, Dataset, DatasetDict, Features, Value, ClassLabel
from PIL import Image


logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

logger = logging.getLogger(__name__)

AUDIO_DIR = 'audio'
VIDEO_DIR = 'video'
IMAGE_DIR = 'image'
EXTENSION_READERS = {
    '.ipynb': lambda f: json.load(f),
    '.md': lambda f: f.read(),
    '.py': lambda f: f.read(),
    '.csv': lambda f: pd.read_csv(f),
    '.json': lambda f: json.load(f),
    '.yaml': lambda f: yaml.safe_load(f),
    '.txt': lambda f: f.read(),
    '.xml': lambda f: f.read(),
    '.html': lambda f: f.read(),
    '.css': lambda f: f.read(),
    '.js': lambda f: f.read(),
    '.h': lambda f: f.read(),
    '.php': lambda f: f.read(),
    '.rb': lambda f: f.read(),
    '.sql': lambda f: f.read(),
    '.xls': lambda f: pd.read_excel(f),
    '.xlsx': lambda f: pd.read_excel(f),
    '.jsp': lambda f: f.read(),
    '.jspx': lambda f: f.read(),
    '.vue': lambda f: f.read(),
    '.ejs': lambda f: f.read(),
    '.erb': lambda f: f.read(),
    '.aspx': lambda f: f.read(),
    '.java': lambda f: f.read(),
    '.cpp': lambda f: f.read(),
    '.h': lambda f: f.read(),
    '.c': lambda f: f.read(),
    '.cc': lambda f: f.read(),
    '.cxx': lambda f: f.read(),
    '.m': lambda f: f.read(),
    '.mm': lambda f: f.read(),
    '.swift': lambda f: f.read(),
    '.objcpp': lambda f: f.read(),
    '.cs': lambda f: f.read(),
    '.go': lambda f: f.read(),
    '.rs': lambda f: f.read(),
    '.kt': lambda f: f.read(),
    '.dart': lambda f: f.read(),
    '.pl': lambda f: f.read(),
    '.pm': lambda f: f.read(),
    '.sh': lambda f: f.read(),
    '.sqlite': lambda f: f.read(),
    '.ppt': lambda f: read_pptx(f),
    '.pptx': lambda f: read_pptx(f)
}

EXTENSION_PATTERN = r".*\.(md|py|csv|json|yaml|txt|xml|html|css|js|java|cpp|h|php|rb|sql|xls|xlsx|ppt|pptx|ipynb|jsp|jspx|vue|ejs|erb|aspx|c|cc|cxx|m|mm|swift|objcpp|cs|go|rs|kt|dart|pl|pm|sh|sqlite)$"
SUPPORTED_EXTENSIONS_IMAGE = [".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".svg", ".png"]
SUPPORTED_EXTENSIONS_AUDIO = [".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma", ".m4a"]
SUPPORTED_EXTENSIONS_VIDEO = [".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv", ".webm"]

def load_and_label_files(root_dir: str) -> DatasetDict:
    """
    Load and label files from a folder and its subfolders.

    Args:
        root_dir (str): The root directory to search for files.

    Returns:
        DatasetDict: A dictionary of datasets containing the labeled files.
        
    Example:
       def main():
           root_dir = "C:/Users/heman/"
           
           try:
               labeled_datasets = load_and_label_files(root_dir)
               print(labeled_datasets['audio']['path'][0] )
               print("Labeled datasets:")
               for file_type, dataset in labeled_datasets.items():
                   print(f"{file_type}: {len(dataset)} files")
           except FileNotFoundError as e:
               print(f"Error: {str(e)}")
           except Exception as e:
               print(f"An unexpected error occurred: {str(e)}")

        if __name__ == "__main__":
           main()
    """
    file_paths: Dict[str, List[str]] = {
        "audio": [],
        "video": [],
        "text": [],
        "image": [],
        "code": []
    }

    audio_extensions = [".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma", ".m4a"]
    video_extensions = [".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv", ".webm"]
    text_extensions = [".txt", ".md", ".rtf", ".doc", ".docx", ".pdf", ".epub"]
    image_extensions = [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".svg", ".webp"]
    code_extensions = [
        ".py", ".java", ".c", ".cpp", ".js", ".ts", ".html", ".css", ".php", ".rb", 
        ".go", ".rs", ".swift", ".kt", ".m", ".sh", ".bat", ".pl", ".lua", ".sql", 
        ".r", ".ipynb", ".json", ".xml", ".yaml", ".yml"
    ]

    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            file_path = os.path.join(dirpath, filename)
            file_path = file_path.replace("\\", "/") 
            extension = os.path.splitext(filename)[1].lower()

            if extension in audio_extensions:
                file_paths["audio"].append(file_path)
            elif extension in video_extensions:
                file_paths["video"].append(file_path)
            elif extension in text_extensions:
                file_paths["text"].append(file_path)
            elif extension in image_extensions:
                file_paths["image"].append(file_path)
            elif extension in code_extensions:
                file_paths["code"].append(file_path)

    datasets = {}
    for file_type, paths in file_paths.items():
        if paths:
            dataset = Dataset.from_dict({"path": paths})
            dataset = dataset.map(lambda x: {"label": file_type}, desc=file_type)
            datasets[file_type] = dataset

    return DatasetDict(datasets)


def is_supported_extension(filename: str,extensions: List[str]) -> bool:
    """
    
    Check if the file has one of the supported extensions.

    Args:
        filename (str): The name of the file.

    Returns:
        bool: True if the file is supported, otherwise False.
    
    
    """
    return any(filename.lower().endswith(ext) for ext in extensions)

def convert_to_flac(input_file: str, output_dir: str) -> Optional[str]: 
    """Convert an audio file to FLAC format using FFmpeg."""
    if not is_supported_extension(input_file, extensions=SUPPORTED_EXTENSIONS_AUDIO):
        logger.error(f"Unsupported file extension for file: {input_file}")
        return None

    output_file = os.path.join(output_dir, os.path.basename(input_file)).replace(os.path.splitext(input_file)[1], '.flac')

    try:
        logger.info(f"Converting {input_file} to {output_file} using FFmpeg")
        subprocess.run(['ffmpeg', '-i', input_file, output_file], check=True)
        logger.info(f"Successfully converted {input_file} to {output_file}")
        return output_file
    except subprocess.CalledProcessError as e:
        logger.error(f"Failed to convert {input_file} to FLAC format. Error: {e}")
        return None
def convert_to_mp4(input_file: str, output_dir: str) -> Optional[str]:
    """Convert a video file to MP4 format using FFmpeg."""
    if not is_supported_extension(input_file, extensions=SUPPORTED_EXTENSIONS_VIDEO):
        logger.error(f"Unsupported file extension for file: {input_file}")
        return None

    output_file = os.path.join(output_dir, os.path.basename(input_file)).replace(os.path.splitext(input_file)[1], '.mp4')

    try:
        logger.info(f"Converting {input_file} to {output_file} using FFmpeg")
        subprocess.run(['ffmpeg', '-i', input_file, output_file], check=True)
        logger.info(f"Successfully converted {input_file} to {output_file}")
        return output_file
    except subprocess.CalledProcessError as e:
        logger.error(f"Failed to convert {input_file} to MP4 format. Error: {e}")
        return None
def convert_files_in_directory_audio(directory: str) -> List[str]:
    """
    Converts all supported audio files in the given directory to .flac format.
    
    Args:
    directory (str): Path to the directory containing audio files.
    
    Returns:
    List[str]: List of paths to the converted .flac files.
    

    """
    if not os.path.isdir(directory):
        logging.error(f"Directory does not exist: {directory}")
        raise NotADirectoryError(f"Directory does not exist: {directory}")

    converted_files = []
    
    for file_name in os.listdir(directory):
        file_path = os.path.join(directory, file_name)
        _, file_extension = os.path.splitext(file_path)
        
        if file_extension in SUPPORTED_EXTENSIONS_AUDIO:
            try:
                converted_file = convert_to_flac(file_path)
                converted_files.append(converted_file)
            except (FileNotFoundError, ValueError, subprocess.CalledProcessError):
                continue
    
    return converted_files




def convert_files_in_directory_video(directory: str) -> List[str]:
    """
    Converts all supported video files in the given directory to .mp4 format.
    
    Args:
    directory (str): Path to the directory containing video files.
    
    Returns:
    List[str]: List of paths to the converted .mp4 files.
    if __name__ == "__main__":
    if len(sys.argv) != 2:
        logging.error("Usage: python script.py /path/to/directory")
        sys.exit(1)
    
    directory_path = sys.argv[1]
    
    try:
        converted_files = convert_files_in_directory(directory_path)
        logging.info(f"Conversion complete. Converted files: {converted_files}")
    except (NotADirectoryError, FileNotFoundError, ValueError) as e:
        logging.error(f"Error: {e}")
        sys.exit(1)

    """
    if not os.path.isdir(directory):
        logging.error(f"Directory does not exist: {directory}")
        raise NotADirectoryError(f"Directory does not exist: {directory}")

    converted_files = []
    
    for file_name in os.listdir(directory):
        file_path = os.path.join(directory, file_name)
        _, file_extension = os.path.splitext(file_path)
        
        if file_extension in SUPPORTED_EXTENSIONS_VIDEO:
            try:
                converted_file = convert_to_mp4(file_path)
                converted_files.append(converted_file)
            except (FileNotFoundError, ValueError, subprocess.CalledProcessError):
                continue
    
    return converted_files


def convert_to_png(input_file: str, output_dir: str) -> Optional[str]: 
    """Convert an image file to PNG format using Pillow."""
    if not is_supported_extension(input_file, extensions=SUPPORTED_EXTENSIONS_IMAGE):
        logger.error(f"Unsupported file extension for file: {input_file}")
        return None

    output_file = os.path.join(output_dir, os.path.basename(input_file)).replace(os.path.splitext(input_file)[1], '.png')

    try:
        logging.info(f"Converting {input_file} to {output_file}")
        with Image.open(input_file) as img:
            img.save(output_file, "PNG")
        logging.info(f"Successfully converted {input_file} to {output_file}")
        return output_file
    except FileNotFoundError:
        logging.error(f"File not found: {input_file}")
        return None
    except Exception as e:
        logger.error(f"Failed to convert {input_file} to PNG format. Error: {e}")
        return None

def convert_files_in_directory_image(directory: str) -> List[str]:
    """
    Converts all supported images files in the given directory to .png format.
    
    Args:
    directory (str): Path to the directory containing images files.
    
    Returns:
    List[str]: List of paths to the converted .png files.
    

    """
    if not os.path.isdir(directory):
        logging.error(f"Directory does not exist: {directory}")
        raise NotADirectoryError(f"Directory does not exist: {directory}")

    converted_files = []
    
    for file_name in os.listdir(directory):
        file_path = os.path.join(directory, file_name)
        _, file_extension = os.path.splitext(file_path)
        
        if file_extension in SUPPORTED_EXTENSIONS_IMAGE:
            try:
                converted_file = convert_to_png(file_path)
                converted_files.append(converted_file)
            except (FileNotFoundError, ValueError, subprocess.CalledProcessError):
                continue
    
    return converted_files


def find_files(directory: str, extensions: List[str]) -> List[str]:
    """
    Recursively finds all files in the given directory that match the given file extensions.
    """
    matches = []
    try:
        for root, _, files in os.walk(directory):
            for file in files:
                if any(file.lower().endswith(ext) for ext in extensions):
                    matches.append(os.path.join(root, file))
    except Exception as e:
        logging.error(f"Error while finding files: {e}")
    return matches

def move_file_with_timestamp(src_path: str, dest_path: str):
    """
    Moves a file to the destination path, appending a timestamp if a file with the same name exists.
    """
    try:
        if os.path.exists(dest_path):
            base_name, file_extension = os.path.splitext(dest_path)
            timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
            dest_path = f"{base_name}_{timestamp}{file_extension}"
        shutil.move(src_path, dest_path)
    except Exception as e:
        logging.error(f"Error while moving file '{src_path}' to '{dest_path}': {e}")

def move_files(files: List[str], output_dir: str):
    """
    Moves the given files into the specified output directory.
    If a file with the same name already exists, it appends a timestamp to the filename.
    """
    for file in files:
        try:
            os.makedirs(output_dir, exist_ok=True)
            file_extension = file.split('.')[-1]
            dest_dir = os.path.join(output_dir, file_extension)
            os.makedirs(dest_dir, exist_ok=True)
            
            base_name = os.path.basename(file)
            dest_file_path = os.path.join(dest_dir, base_name)
            
            move_file_with_timestamp(file, dest_file_path)
        except Exception as e:
            logging.error(f"Error while processing file '{file}': {e}")

def organize_files(input_dir: str, output_dir: str,extensions: List[str]):
    """
    Organizes files by their extension from the input directory into the output directory.
    organize_files('/content/drive/MyDrive', '/content/hemanth')
    extensions = ['.pdf', '.txt', '.md', '.docx',  '.csv', ]
    """
    try:
        if not os.path.isdir(input_dir):
            logging.error(f"Input directory '{input_dir}' does not exist or is not a directory.")
            return
        if  os.path.isdir(output_dir):
            logging.info(f"Output directory '{output_dir}' does not exist or is not a directory.")
            os.makedirs(output_dir, exist_ok=True)
            logging.info(f"Output directory '{output_dir}' we forcefull creating")
        
        else:
            logging.info(f"Output directory '{output_dir}' Automatic created .")        
            files_to_move = find_files(input_dir, extensions)
            move_files(files_to_move, output_dir)
            logging.info(f"completed moving files into folders {output_dir} 📁 in organized")
            return
    except Exception as e:
        logging.error(f"Error while organizing files: {e}")


def read_pptx(file):
    """Custom function to read .pptx files with python-pptx"""
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)



def get_files_with_extensions(dir_path: str) -> Dict[str, List[str]]:
    ext_files = defaultdict(list)
    for root, dirs, files in os.walk(dir_path):
        for file in files:
            file_path = os.path.join(root, file).replace("\\", "/")
            _, ext = os.path.splitext(file)
            ext_files[ext].append(file_path)
    return ext_files


def write_to_csv(file_path: str, ext_files: Dict[str, List[str]]) -> None:
    with open(file_path, 'w', encoding='utf-8',newline='') as csvfile:
        writer = csv.writer(csvfile)
        max_len = max(len(v) for v in ext_files.values())
        writer.writerow(ext_files.keys())
        for i in range(max_len):
            row = [ext_files[k][i] if i < len(ext_files[k]) else '' for k in ext_files.keys()]
            writer.writerow(row)


def clean_text(text: str) -> Optional[str]:
    """
    Clean the extracted text from the Document using advanced techniques.
    This function can be customized based on the cleaning requirements.

    Parameters:
    - text (str): The text extracted from the PDF.

    Returns:
    - Optional[str]: The cleaned text, or None if the input is invalid.
    """
    if not isinstance(text, str):
        return None
    cleaned_text = text.strip()
    cleaned_text = cleaned_text.lower()
    if not cleaned_text:
        return None
    return cleaned_text

def split_and_save_text(cleaned_text: str, base_output_path: Path, max_size_bytes: int = 50 * 1024 * 1024) -> None:
    """
    Split the cleaned text into multiple files, each smaller than the specified max size, and save them.

    Parameters:
    - cleaned_text (str): The cleaned text to be split and saved.
    - base_output_path (Path): The base path where the text files will be saved.
    - max_size_bytes (int): Maximum size of the text file in bytes.
    """
    part_num = 1
    text_part = ""
    for line in cleaned_text.split('\n'):
        if len(text_part.encode('utf-8')) + len(line.encode('utf-8')) < max_size_bytes:
            text_part += line + '\n'
        else:
            output_path = base_output_path.with_suffix(f'.part{part_num}.txt')
            with open(output_path, 'w', encoding='utf-8') as file:
                file.write(text_part)
            part_num += 1
            text_part = line + '\n' 
    if text_part:
        output_path = base_output_path.with_suffix(f'.part{part_num}.txt')
        with open(output_path, 'w', encoding='utf-8') as file:
            file.write(text_part)


def convert_pdf_to_text(pdf_path: str, output_folder: str) -> None:
    """
    Convert a PDF file to text files, splitting contents to ensure each resulting file is less than 50 MB.

    Parameters:
    - pdf_path (str): Path to the PDF file.
    - output_folder (str): Path to the folder where the text files will be saved.
    """
    output_folder_path = Path(output_folder)
    output_folder_path.mkdir(parents=True, exist_ok=True)
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        cleaned_text = clean_text(text)
        base_output_path = output_folder_path / Path(pdf_path).stem
        split_and_save_text(cleaned_text, base_output_path)
    except Exception as e:
        return None


def process_pdfs_from_csv(csv_path: str, output_folder: str,extension:str='.pdf') -> None:
    """
    Process PDFs listed in a CSV file, converting them to text files and ensuring each part is less than 50 MB.

    Parameters:
    - csv_path (str): Path to the CSV file containing paths to PDF files.
    - output_folder (str): Path to the folder where text files will be stored.
    """
    encodings = ['utf-8', 'latin1', 'utf-16', 'utf-32','ascii']

    for encoding in encodings:
        try:
            pdf_paths = pd.read_csv(csv_path, encoding=encoding)
            break 
        except UnicodeDecodeError:
            continue  
    else:
        raise ValueError("Unable to decode the CSV file with the provided encodings.")

    for pdf_path in pdf_paths[extension]:
        convert_pdf_to_text(pdf_path, output_folder)

def list_files_with_extensions(directory_path):
    try:
        files = os.listdir(directory_path)
        return [file for file in files if re.match(EXTENSION_PATTERN, file)]
    except FileNotFoundError:
        return None

def read_file_content(directory_path, filename):
    try:
        extension = os.path.splitext(filename)[1]
        with open(os.path.join(directory_path, filename), 'r',encoding='utf-8') as file:
            file_reader = EXTENSION_READERS.get(extension)
            return file_reader(file) if file_reader else None
    except Exception as e:
        return None



def process_files_txtfile(directory_path: str, user_folder: str, txt_file_counter: int = 1, txt_file_size: int = 0,size_file:int=5*1024*1024) -> Optional[str]:
    """
    This function recursively processes all files in a given directory and its subdirectories,
    and writes their content to a user-specific text file. Each text file is ensured to be less than 50 MB in size.

    Args:
    directory_path (str): The path to the directory containing the files to be processed.
    user_folder (str): The name of the user-specific folder where the text files will be written.
    txt_file_counter (int): The counter for the current text file.
    txt_file_size (int): The current size of the text file.

    Returns:
    str: The path to the user-specific folder, or None if an error occurred.
    """
    
    user_folder_path = os.path.join(directory_path, user_folder)
    os.makedirs(user_folder_path, exist_ok=True)
    files = list_files_with_extensions(directory_path)
    if files is None:
        return
    for filename in files:
        content = read_file_content(directory_path, filename)
        if content is not None:          
            if txt_file_size >= size_file:
                txt_file_counter += 1
                txt_file_size = 0
            txt_file_path = os.path.join(user_folder_path, f"{txt_file_counter}.txt")
            with open( txt_file_path ,"a", encoding='utf-8') as f:
                if isinstance(content, pd.DataFrame):
                    content_csv = content.to_csv(index=False)
                    f.write(content_csv)
                    txt_file_size += len(content_csv)
                elif isinstance(content, dict):
                    content_str = json.dumps(content, indent=4)
                    f.write(content_str)
                    txt_file_size += len(content_str)
                else:
                    content_str = str(content)
                    f.write(content_str)
                    txt_file_size += len(content_str)
    for subdir in [d for d in os.listdir(directory_path) if os.path.isdir(os.path.join(directory_path, d))]:
        subdir_path = os.path.join(directory_path, subdir)
        txt_file_counter, txt_file_size = process_files_txtfile(subdir_path, user_folder, txt_file_counter, txt_file_size)
    return txt_file_counter, txt_file_size

def reformat_txt_files(folders: List[str]) -> None:
    """
    Reformats .txt files in the specified folders to give them a clean, professional look.
    Saves the reformatted files in new folders while keeping the original files intact.

    Args:
        folders (List[str]): List of folder paths containing the .txt files to reformat.

    Returns:
        None
    """
    for folder in folders:
        new_folder = os.path.join(folder, "reformatted")
        os.makedirs(new_folder, exist_ok=True)

        for filename in os.listdir(folder):
            if filename.endswith(".txt"):
                file_path = os.path.join(folder, filename)
                with open(file_path, "r",encoding='utf-8') as file:
                    content = file.read()
                content = content.strip()
                content = " ".join(content.split())  
                content = content.replace("\n", "\n\n")  
                lines = content.split("\n")
                if len(lines) > 50:
                    content = ""
                    for i, line in enumerate(lines, start=1):
                        content += line + "\n"
                        if i % 50 == 0:
                            content += f"\nPage {i // 50}\n\n"
                new_filename = "_".join(filename.split()) 
                new_file_path = os.path.join(new_folder, new_filename)
                with open(new_file_path, "w", encoding='utf-8') as file:
                    file.write(content)




def convert_files_in_directory_text(dir_path:str,output_dir:str,csv_path:str):
    logging.info(f"extract all data from {dir_path} to {output_dir}")
    ext_files = get_files_with_extensions(dir_path)
    write_to_csv(csv_path, ext_files)
    process_pdfs_from_csv(csv_path=csv_path, output_folder=output_dir, extension='.pdf')
    process_files_txtfile(dir_path,  output_dir)
    reformat_txt_files([output_dir])
    for filename in os.listdir(output_dir):
       if filename.endswith(".txt"):
         file_path = os.path.join(output_dir, filename)
         try:
           os.remove(file_path)
         except OSError as e:
           print(f"Error deleting file '{file_path}': {e}")
    logging.info(f"completely extracted data into {output_dir}")


def create_output_dirs(base_output_dir: str) -> Dict[str, str]:
    """Create output directories for different file types."""
    dirs = {}
    for dir_name in [AUDIO_DIR, VIDEO_DIR, IMAGE_DIR]:
        path = os.path.join(base_output_dir, dir_name)
        os.makedirs(path, exist_ok=True)
        dirs[dir_name] = path
    return dirs


def convert_file(file_path: str, output_dir: str,
                 convert_function: callable) -> None:
    """Convert a single file and log the result."""
    try:
        converted_file_path = convert_function(file_path, output_dir) 
        if converted_file_path:
            logger.info(f"Successfully converted: {file_path} to {converted_file_path}")
        else:
            logger.warning(f"Conversion skipped for: {file_path}")
    except Exception as e:
        logger.error(f"Failed to convert {file_path}: {e}")


def convert_files_in_directory(directory: str, output_dir: str) -> None:
    """Convert and organize files in a directory by type."""
    directory = os.path.normpath(directory) + os.sep
    output_dir = os.path.normpath(output_dir) + os.sep

    if not os.path.isdir(directory):
        logger.error(f"Directory does not exist: {directory}")
        return

    output_dirs = create_output_dirs(output_dir)
    output_dirs_set = set(output_dirs.values())

    for root, _, files in os.walk(directory):
        root = os.path.normpath(root) + os.sep
        # Skip output directories to prevent infinite recursion
        if root in output_dirs_set:
            continue

        for file_name in files:
            file_path = os.path.join(root, file_name)

            if is_supported_extension(file_name, SUPPORTED_EXTENSIONS_AUDIO):
                convert_file(file_path, output_dirs[AUDIO_DIR], convert_to_flac)
            elif is_supported_extension(file_name, SUPPORTED_EXTENSIONS_VIDEO):
                convert_file(file_path, output_dirs[VIDEO_DIR], convert_to_mp4)
            elif is_supported_extension(file_name, SUPPORTED_EXTENSIONS_IMAGE):
                convert_file(file_path, output_dirs[IMAGE_DIR], convert_to_png)


