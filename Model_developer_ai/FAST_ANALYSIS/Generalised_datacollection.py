import os
import logging
from typing import List, Optional
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor
import shutil
import pandas as pd
import json
import yaml
import docx
import pptx
import fitz  # PyMuPDF


logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    filename='file_processor.log',
    filemode='w'
)
logger = logging.getLogger(__name__)

def read_file_content(file_path: Path) -> Optional[str]:
    """
    Read content from various file types and return as string.

    Args:
        file_path (Path): Path to the file.

    Returns:
        Optional[str]: Content of the file as a string, or None if unsuccessful.
    """
    ext = file_path.suffix.lower()
    try:
        if ext in ['.txt', '.py', '.java', '.cpp', '.html', '.css', '.js', '.xml']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext in ['.csv', '.xlsx', '.xls']:
            df = pd.read_excel(file_path) if ext != '.csv' else pd.read_csv(file_path)
            return df.to_string()
        elif ext == '.json':
            with open(file_path, 'r') as f:
                return json.dumps(json.load(f), indent=2)
        elif ext in ['.yaml', '.yml']:
            with open(file_path, 'r') as f:
                return yaml.dump(yaml.safe_load(f))
        elif ext == '.docx':
            doc = docx.Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif ext in ['.ppt', '.pptx']:
            prs = pptx.Presentation(file_path)
            return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
        elif ext == '.pdf':
            with fitz.open(file_path) as doc:
                return '\n'.join([page.get_text() for page in doc])
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return None
    except Exception as e:
        logger.error(f"Error reading file {file_path}: {str(e)}")
        return None

def process_file(file_path: Path, output_folder: Path, max_size: int = 10 * 1024 * 1024) -> None:
    """
    Process a single file and convert it to .txt if size is less than max_size.

    Args:
        file_path (Path): Path to the input file.
        output_folder (Path): Path to the output folder.
        max_size (int): Maximum file size in bytes (default: 10MB).
    """
    try:
        if file_path.stat().st_size < max_size:
            content = read_file_content(file_path)
            if content is not None:
                dest_path = output_folder / (file_path.stem + '.txt')
                dest_path.parent.mkdir(parents=True, exist_ok=True)
                with open(dest_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                logger.info(f"Converted and copied file: {file_path} to {dest_path}")
            else:
                logger.warning(f"Skipped file (unable to read content): {file_path}")
        else:
            logger.warning(f"Skipped file (size > {max_size} bytes): {file_path}")
    except Exception as e:
        logger.error(f"Error processing file {file_path}: {str(e)}")

def process_folder(folder_path: Path, output_folder: Path, max_size: int = 10 * 1024 * 1024) -> None:
    """
    Recursively process a folder and convert files to .txt with size less than max_size.

    Args:
        folder_path (Path): Path to the input folder.
        output_folder (Path): Path to the output folder.
        max_size (int): Maximum file size in bytes (default: 10MB).
    """
    try:
        for item in folder_path.rglob('*'):
            if item.is_file():
                process_file(item, output_folder, max_size)
    except Exception as e:
        logger.error(f"Error processing folder {folder_path}: {str(e)}")

def convertfolders_to_txtfolder(input_folders: List[str], output_folder: str, max_workers: int = 4) -> None:
    """
    Process multiple input folders concurrently.

    Args:
        input_folders (List[str]): List of input folder paths.
        output_folder (str): Path to the output folder.
        max_workers (int): Maximum number of concurrent workers.
    """
    output_path = Path(output_folder)
    output_path.mkdir(parents=True, exist_ok=True)

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = []
        for folder in input_folders:
            folder_path = Path(folder)
            if folder_path.is_dir():
                futures.append(executor.submit(process_folder, folder_path, output_path))
            else:
                logger.warning(f"Skipping invalid folder: {folder}")

        for future in futures:
            future.result()